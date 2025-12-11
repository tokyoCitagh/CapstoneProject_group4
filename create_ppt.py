import importlib
import collections
# Some Python builds expose `collections.abc` as a submodule but not as an attribute
# on the `collections` module. Ensure the attribute exists so third-party libs
# (e.g. python-pptx compatibility helpers) that access `collections.abc` via
# attribute lookup don't crash with AttributeError.
if not hasattr(collections, 'abc'):
    collections.abc = importlib.import_module('collections.abc')

from pptx import Presentation
from pptx.util import Inches, Pt
import os

out_name = "EcomApp_Project_Documentation.pptx"
out_path = os.path.join(os.path.dirname(__file__), out_name)

prs = Presentation()

def add_title_slide(title, subtitle=None):
    sld = prs.slides.add_slide(prs.slide_layouts[0])
    sld.shapes.title.text = title
    if subtitle:
        try:
            sld.placeholders[1].text = subtitle
        except Exception:
            pass


def add_bullet_slide(title, bullets):
    sld = prs.slides.add_slide(prs.slide_layouts[1])
    sld.shapes.title.text = title
    body = sld.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0


add_title_slide("Image Electronics — EcomApp", "Project summary & deployment notes")

add_bullet_slide("Project Overview", [
    "Django-based e-commerce site (store, services, admin/portal)",
    "Features: checkout, shipping (pickup/delivery), analytics, staff portal",
    "Deployment: Railway with PostgreSQL, Cloudinary for media",
])

add_bullet_slide("Key Features Implemented", [
    "Checkout: phone persistence, shipping options (pickup/delivery, speeds)",
    "TomTom-based distance shipping estimate endpoint",
    "Staff portal: order fulfillment, shipping speed, colored status badges",
    "Traffic analytics: unique sessions dataset, Traffic Trends UI"
])

add_bullet_slide("Important Files & Modules", [
    "`store/views.py` — checkout, portal views, shipping estimate",
    "`store/models.py` — Customer.phone, ShippingAddress.phone, Order.shipping_fee/fulfillment",
    "Templates: `store/orders_list.html`, `store/order_detail.html`",
    "Migrations: `store/migrations/0014` and `0015` for phone & shipping fields",
])

add_bullet_slide("Database & Migrations", [
    "Used Django migrations for schema changes (applied in production)",
    "Temporary secure migration endpoint was used once to run migrations safely",
    "Defensive code paths added then removed after migrations applied",
])

add_bullet_slide("Deployment & Domain", [
    "Hosted on Railway (web service) with Gunicorn + WhiteNoise",
    "Custom domain `www.imageelectronics.org` added and TLS via Let's Encrypt",
    "ALLOWED_HOSTS and CSRF_TRUSTED_ORIGINS updated to accept custom domain",
])

add_bullet_slide("Troubleshooting Notes", [
    "HTTP 400 due to DisallowedHost — fixed by adding domain to ALLOWED_HOSTS",
    "CSRF 403 on login — fixed by adding domain to CSRF_TRUSTED_ORIGINS and redeploy",
    "DNS apex redirect (imageelectronics.org) recommended: URL redirect to www",
])

add_bullet_slide("How to Run Locally", [
    "Create virtualenv: `python -m venv env`",
    "Install: `pip install -r requirements.txt`",
    "Run migrations: `python manage.py migrate`",
    "Start dev server: `python manage.py runserver`",
])

add_bullet_slide("Next Steps / Recommendations", [
    "Harden cookies: set SESSION_COOKIE_SECURE and CSRF_COOKIE_SECURE in prod",
    "Consider server-side shipping fee calculation to prevent client tampering",
    "Add end-to-end tests for checkout and portal flows",
])

add_bullet_slide("Contacts & Notes", [
    "Repository: CapstoneProject_group4 (owner: tokyoCitagh)",
    "Deployment: Railway project `profound-surprise` (production environment)",
    "Created by: development session — see repo for commit history",
])

prs.save(out_path)
print(f"Saved PowerPoint to: {out_path}")
